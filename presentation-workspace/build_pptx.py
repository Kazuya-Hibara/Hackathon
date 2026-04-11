"""Build Lazy Captain presentation PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS = os.path.join(BASE, "screenshots")
IMAGES = os.path.join(BASE, "images")
OUTPUT = os.path.join(BASE, "lazy-captain-presentation.pptx")

# Colors
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0xAA)
GREEN = RGBColor(0x34, 0xD3, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=BG_DARK):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, font_size=18, color=WHITE, bold=False, space_before=Pt(8)):
    """Add paragraph to existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.space_before = space_before
    return p


def add_image(slide, path, left, top, width=None, height=None):
    """Add image to slide."""
    kwargs = {"left": Inches(left), "top": Inches(top)}
    if width:
        kwargs["width"] = Inches(width)
    if height:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(path, **kwargs)


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Title
add_text(slide, 1.5, 1.5, 10, 1.2, "Lazy Captain", font_size=54, color=CYAN, bold=True)
add_text(slide, 1.5, 2.8, 10, 0.8, "Your AI copilot for managerial impact", font_size=28, color=GRAY)

# Screenshot on right
add_image(slide, os.path.join(SCREENSHOTS, "today-tab.png"), 6.5, 3.2, width=6)

# Subtle tagline
add_text(slide, 1.5, 5.5, 5, 0.5, "TEAMZ AI Hackathon 2026  |  Tokyo", font_size=14, color=GRAY)


# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 12, 1, "Your best work has no proof", font_size=42, color=WHITE, bold=True)

# Problem illustration on right
add_image(slide, os.path.join(IMAGES, "problem.png"), 7, 1.5, width=5.5)

# 3 bullet points - left side
tf = add_text(slide, 1, 2.2, 5.5, 4, "", font_size=22, color=GRAY)
tf.paragraphs[0].text = ""

bullets = [
    "You unblocked 3 teams today.\nWhere's the receipt?",
    "Decisions, coaching, alignment\n  none leave an artifact.",
    "Review time comes around:\n\"What did I even do last quarter?\"",
]
for b in bullets:
    p = add_para(tf, b, font_size=20, color=GRAY, space_before=Pt(24))


# ============================================================
# SLIDE 3: Meet Lazy Captain
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 12, 1, "One place for everything you did", font_size=42, color=WHITE, bold=True)

# Full-page screenshot
add_image(slide, os.path.join(SCREENSHOTS, "today-full.png"), 7, 0.8, height=6)

# Key points on left
tf = add_text(slide, 1, 2, 5.5, 4.5, "", font_size=20)
tf.paragraphs[0].text = ""

points = [
    ("AI reads your tools", "Calendar, Slack, web activity"),
    ("Suggests impact items", "You just approve, edit, or dismiss"),
    ("7 categories", "Decision  |  Blocker Removed  |  Alignment\nStructure  |  Team Dev  |  Strategic  |  Fires"),
]
for title, desc in points:
    add_para(tf, title, font_size=22, color=CYAN, bold=True, space_before=Pt(20))
    add_para(tf, desc, font_size=16, color=GRAY, space_before=Pt(4))


# ============================================================
# SLIDE 4: The Daily Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.4, 12, 1, "5 minutes. Done.", font_size=42, color=WHITE, bold=True)
add_text(slide, 0.8, 1.3, 12, 0.6, "The daily workflow that runs itself", font_size=20, color=GRAY)

# Flow diagram
add_image(slide, os.path.join(IMAGES, "flow.png"), 1, 2.2, width=11)

# Brief description below
add_text(slide, 1, 5.5, 11, 1.5,
         "Connectors pull your data  >  AI analyzes patterns  >  You review suggestions  >  Impact log builds  >  Day summary auto-generated",
         font_size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: It Learns You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 12, 1, "It learns you", font_size=42, color=WHITE, bold=True)

# Screenshot showing suggestion cards
add_image(slide, os.path.join(SCREENSHOTS, "today-tab.png"), 6.5, 1.5, width=6)

# Left content
tf = add_text(slide, 1, 2, 5, 4.5, "", font_size=20)
tf.paragraphs[0].text = ""

learn_points = [
    ("Accept", "AI nailed it. One tap."),
    ("Edit", "Close but not quite. Fix and save."),
    ("Dismiss", "Not relevant. AI learns to skip next time."),
]
for action, desc in learn_points:
    add_para(tf, action, font_size=24, color=GREEN, bold=True, space_before=Pt(24))
    add_para(tf, desc, font_size=16, color=GRAY, space_before=Pt(4))

add_para(tf, "", font_size=12, space_before=Pt(20))
add_para(tf, "Mem9 stores your preferences.", font_size=18, color=PURPLE, space_before=Pt(4))
add_para(tf, "Suggestions sharpen over time.", font_size=18, color=PURPLE, space_before=Pt(4))


# ============================================================
# SLIDE 6: Growth Dashboard
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 12, 1, "See your patterns", font_size=42, color=WHITE, bold=True)

# Growth tab screenshot (hero)
add_image(slide, os.path.join(SCREENSHOTS, "growth-tab.png"), 1.5, 1.5, width=10)

# Caption bar at bottom
tf = add_text(slide, 1, 6.2, 11, 1, "", font_size=16, color=GRAY, align=PP_ALIGN.CENTER)
tf.paragraphs[0].text = "Category distribution  |  Impact trends  |  Scope tracking  |  Delegation scores"


# ============================================================
# SLIDE 7: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.4, 12, 1, "Built with 5 sponsor technologies", font_size=36, color=WHITE, bold=True)

# Architecture diagram
add_image(slide, os.path.join(IMAGES, "architecture.png"), 1, 1.5, width=7)

# Tech list on right
tf = add_text(slide, 8.5, 1.8, 4.5, 5, "", font_size=18)
tf.paragraphs[0].text = ""

techs = [
    ("Agnes AI", "LLM suggestion engine"),
    ("TiDB Cloud Zero", "Serverless MySQL database"),
    ("Bright Data", "SERP enrichment API"),
    ("Mem9", "Persistent learning memory"),
    ("Zeabur", "Cloud deployment platform"),
]
for name, desc in techs:
    add_para(tf, name, font_size=20, color=CYAN, bold=True, space_before=Pt(16))
    add_para(tf, desc, font_size=14, color=GRAY, space_before=Pt(2))

add_para(tf, "", font_size=10, space_before=Pt(24))
add_para(tf, "FastAPI + Vanilla JS + Python", font_size=14, color=GRAY, space_before=Pt(4))


# ============================================================
# SLIDE 8: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 1.5, 2, 10, 1.2, "Lazy Captain", font_size=60, color=CYAN, bold=True,
         align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 3.5, 10, 0.8, "Your AI copilot for managerial impact",
         font_size=24, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 5, 10, 0.6, "Built at TEAMZ AI Hackathon 2026  |  Tokyo",
         font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Sponsor names row
add_text(slide, 1, 6, 11, 0.5,
         "Agnes AI   |   TiDB   |   Bright Data   |   Mem9   |   Zeabur",
         font_size=16, color=PURPLE, align=PP_ALIGN.CENTER)


# ============================================================
# Save
# ============================================================
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
